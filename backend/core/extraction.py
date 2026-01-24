"""
BeePrepared - File Extraction Service
Converts all supported file formats to raw text for processing.

Supported Formats:
- PDF: PyMuPDF (fitz) extraction
- PPTX: python-pptx extraction
- MD/TXT: Direct read
- Audio (WAV/MP3): Azure Speech transcription
- Video (MP4/etc): FFmpeg audio extraction → Azure Speech
"""

import os
import logging
import tempfile
import time
import shutil
from pathlib import Path
from typing import Optional, Tuple
from dotenv import load_dotenv

import boto3
import azure.cognitiveservices.speech as speechsdk
import fitz
from pptx import Presentation
from docx import Document
import ffmpeg

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

load_dotenv()


class ExtractionService:
    """
    Extracts text content from various file formats.
    """
    
    SUPPORTED_AUDIO = {'.wav', '.mp3', '.m4a', '.flac', '.ogg', '.wma'}
    SUPPORTED_VIDEO = {'.mp4', '.mov', '.avi', '.mkv', '.webm', '.wmv'}
    SUPPORTED_DOCS = {'.pdf', '.pptx', '.ppt', '.md', '.txt', '.docx'}
    
    def __init__(self):
        # self._setup_azure_speech()
        self._setup_deepgram()
        self._setup_r2()
    
    def _setup_deepgram(self):
        """Initialize Deepgram SDK."""
        self.dg_key = os.getenv("DEEPGRAM_API_KEY") or os.getenv("DEEPGRAM_KEY")
        if self.dg_key:
            try:
                from deepgram import DeepgramClient
                self.deepgram = DeepgramClient(api_key=self.dg_key)
                logger.info("Deepgram SDK initialized successfully.")
            except Exception as e:
                logger.error(f"Failed to initialize Deepgram: {e}")
                self.deepgram = None
        else:
            logger.warning("DEEPGRAM_API_KEY not found. Audio transcription will fail.")
            self.deepgram = None

    def _setup_r2(self):
        """Initialize Cloudflare R2 client."""
        self.r2_endpoint = os.getenv("R2_ENDPOINT_URL")
        self.r2_key = os.getenv("R2_ACCESS_KEY_ID")
        self.r2_secret = os.getenv("R2_SECRET_ACCESS_KEY")
        self.r2_bucket = os.getenv("R2_BUCKET_NAME")

        if self.r2_endpoint and self.r2_key and self.r2_secret:
            try:
                self.s3_client = boto3.client(
                    service_name='s3',
                    endpoint_url=self.r2_endpoint,
                    aws_access_key_id=self.r2_key,
                    aws_secret_access_key=self.r2_secret
                )
            except Exception as e:
                logger.error(f"Failed to initialize R2 client: {e}")
                self.s3_client = None
        else:
            logger.warning("R2 credentials missing. R2 download will fail.")
            self.s3_client = None

    def extract_from_r2(self, r2_key: str) -> Tuple[str, dict]:
        """
        Download file from R2 -> Temp File -> Extract -> Return.
        """
        if not self.s3_client:
            raise RuntimeError("R2 client not initialized. Check credentials.")

        logger.info(f"Downloading from R2: {r2_key}")
        ext = os.path.splitext(r2_key)[1]
        
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
            tmp_path = tmp.name
        
        try:
            self.s3_client.download_file(self.r2_bucket, r2_key, tmp_path)
            logger.info(f"Downloaded to temp file: {tmp_path}")
            return self.extract(tmp_path)
            
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    
    def detect_file_type(self, file_path: str) -> str:
        """Detect the type of file based on extension."""
        ext = Path(file_path).suffix.lower()
        
        if ext in self.SUPPORTED_AUDIO:
            return "audio"
        elif ext in self.SUPPORTED_VIDEO:
            return "video"
        elif ext == ".pdf":
            return "pdf"
        elif ext in {".pptx", ".ppt"}:
            return "pptx"
        elif ext == ".md":
            return "markdown"
        elif ext == ".txt":
            return "text"
        elif ext == ".docx":
            return "docx"
        else:
            return "unknown"
    
    def extract(self, file_path: str) -> Tuple[str, dict]:
        """
        Main extraction method. Routes to appropriate handler.
        
        Returns:
            Tuple[str, dict]: (extracted_text, metadata)
        """
        file_type = self.detect_file_type(file_path)
        logger.info(f"Extracting from {file_path} (type: {file_type})")
        
        extractors = {
            "pdf": self._extract_pdf,
            "pptx": self._extract_pptx,
            "markdown": self._extract_text,
            "text": self._extract_text,
            "docx": self._extract_docx,
            "audio": self._transcribe_deepgram,
            "video": self._transcribe_video,
        }
        
        extractor = extractors.get(file_type)
        if not extractor:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        start_time = time.time()
        text, metadata = extractor(file_path)
        elapsed = time.time() - start_time
        
        metadata["file_type"] = file_type
        metadata["extraction_time_seconds"] = round(elapsed, 2)
        metadata["char_count"] = len(text)
        metadata["word_count"] = len(text.split())
        
        logger.info(f"Extraction complete: {metadata['word_count']} words in {elapsed:.2f}s")
        return text, metadata
    
    # =========================================================================
    # DOCUMENT EXTRACTORS
    # =========================================================================
    
    def _extract_pdf(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from PDF using PyMuPDF."""
        logger.info(f"Extracting PDF: {file_path}")
        
        doc = fitz.open(file_path)
        text_parts = []
        
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text("text")
            if page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
        
        doc.close()
        
        full_text = "\n\n".join(text_parts)
        metadata = {
            "page_count": len(text_parts),
            "source": Path(file_path).name
        }
        
        return full_text, metadata
    
    def _extract_pptx(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from PowerPoint presentations."""
        logger.info(f"Extracting PPTX: {file_path}")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"[Speaker Notes: {notes}]")
            
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        
        full_text = "\n\n".join(text_parts)
        metadata = {
            "slide_count": len(prs.slides),
            "source": Path(file_path).name
        }
        
        return full_text, metadata
    
    def _extract_text(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from plain text or markdown files."""
        logger.info(f"Reading text file: {file_path}")
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        
        metadata = {
            "source": Path(file_path).name
        }
        
        return text, metadata
    
    def _extract_docx(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from Word documents."""
        logger.info(f"Extracting DOCX: {file_path}")
        
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        full_text = "\n\n".join(text_parts)
        metadata = {
            "paragraph_count": len(doc.paragraphs),
            "source": Path(file_path).name
        }
        
        return full_text, metadata
    
    # =========================================================================
    # AUDIO/VIDEO TRANSCRIPTION
    # =========================================================================
    
    def _convert_to_wav(self, input_path: str, output_path: str) -> None:
        """Convert audio/video to Azure-compatible WAV format."""
        logger.info(f"Converting to WAV: {input_path} -> {output_path}")
        try:
            (
                ffmpeg
                .input(input_path)
                .output(output_path, acodec='pcm_s16le', ac=1, ar='16000')
                .overwrite_output()
                .run(quiet=True, capture_stdout=True, capture_stderr=True)
            )
        except ffmpeg.Error as e:
            logger.error(f"FFmpeg error: {e.stderr.decode() if e.stderr else str(e)}")
            raise
    
    def _transcribe_audio(self, file_path: str) -> Tuple[str, dict]:
        """Transcribe audio file using Azure Speech."""
        if not self.speech_config:
            raise RuntimeError("Azure Speech not configured. Check AZURE_SPEECH_KEY.")
        
        # Convert to WAV if needed
        ext = Path(file_path).suffix.lower()
        if ext != ".wav":
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                wav_path = tmp.name
            self._convert_to_wav(file_path, wav_path)
        else:
            wav_path = file_path
        
        logger.info(f"Transcribing audio: {wav_path}")
        
        try:
            return self._run_continuous_recognition(wav_path)
        finally:
            # Cleanup temp file if created
            if ext != ".wav" and os.path.exists(wav_path):
                os.remove(wav_path)
    
    def _transcribe_deepgram(self, file_path: str) -> Tuple[str, dict]:
        """
        Transcribe audio using Deepgram Nova-2 (Direct REST API).
        Bypasses SDK version issues.
        """
        if not self.dg_key:
            self.dg_key = os.getenv("DEEPGRAM_API_KEY") or os.getenv("DEEPGRAM_KEY")
            
        if not self.dg_key:
            raise RuntimeError("Deepgram API Key not found.")

        logger.info(f"Transcribing with Deepgram Nova-2 (REST): {file_path}")
        
        try:
            import requests
            
            with open(file_path, "rb") as audio:
                data = audio.read()
            
            # Determine content type (optional, but good practice)
            ext = Path(file_path).suffix.lower()
            content_type = "audio/wav"
            if ext == ".mp3": content_type = "audio/mpeg"
            elif ext == ".m4a": content_type = "audio/mp4" # approx
            
            url = "https://api.deepgram.com/v1/listen?model=nova-2&smart_format=true&punctuate=true&language=en"
            headers = {
                "Authorization": f"Token {self.dg_key}",
                "Content-Type": content_type
            }
            
            response = requests.post(url, headers=headers, data=data)
            response.raise_for_status()
            
            result = response.json()
            # Parse result
            transcript = result['results']['channels'][0]['alternatives'][0]['transcript']
            
            logger.info(f"Deepgram transcription complete. Length: {len(transcript)} chars")
            
            metadata = {
                "source": Path(file_path).name,
                "provider": "deepgram",
                "model": "nova-2 (REST)"
            }
            return transcript, metadata

        except Exception as e:
            logger.error(f"Deepgram transcription failed: {e}")
            raise

    def _transcribe_video(self, file_path: str) -> Tuple[str, dict]:
        """Extract audio from video and transcribe with Deepgram."""
        logger.info(f"Extracting audio from video: {file_path}")
        
        # Extracted audio temp file
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            wav_path = tmp.name
        
        try:
            self._convert_to_wav(file_path, wav_path)
            # Re-route to Deepgram
            return self._transcribe_deepgram(wav_path)
        finally:
            if os.path.exists(wav_path):
                os.remove(wav_path)
    
    def _run_continuous_recognition(self, wav_path: str) -> Tuple[str, dict]:
       """Deprecated Azure Method - Kept for reference but unused."""
       pass


# =============================================================================
# CONVENIENCE FUNCTION
# =============================================================================

def extract_text(file_path: str) -> Tuple[str, dict]:
    """
    Convenience function to extract text from any supported file.
    
    Usage:
        text, metadata = extract_text("/path/to/file.pdf")
    """
    service = ExtractionService()
    return service.extract(file_path)


# =============================================================================
# CLI TEST
# =============================================================================

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python extraction.py <file_path>")
        print("\nSupported formats:")
        print("  Documents: PDF, PPTX, DOCX, MD, TXT")
        print("  Audio: WAV, MP3, M4A, FLAC, OGG")
        print("  Video: MP4, MOV, AVI, MKV, WEBM")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    if not os.path.exists(file_path):
        print(f"Error: File not found: {file_path}")
        sys.exit(1)
    
    print(f"\n{'='*60}")
    print(f"  BeePrepared - Text Extraction")
    print(f"{'='*60}")
    print(f"File: {file_path}")
    
    try:
        service = ExtractionService()
        text, metadata = service.extract(file_path)
        
        print(f"\n{'='*60}")
        print("METADATA:")
        print(f"{'='*60}")
        for key, value in metadata.items():
            print(f"  {key}: {value}")
        
        print(f"\n{'='*60}")
        print("EXTRACTED TEXT (first 2000 chars):")
        print(f"{'='*60}")
        print(text[:2000])
        
        if len(text) > 2000:
            print(f"\n... [{len(text) - 2000} more characters]")
        
        print(f"\n{'='*60}")
        print("✅ Extraction successful!")
        print(f"{'='*60}")
        
    except Exception as e:
        print(f"\n❌ Extraction failed: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
