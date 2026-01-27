import os
import sys
import logging
import json
import time
from pathlib import Path

# Ensure backend path is in sys.path
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from backend.core.ingest import IngestionService
from backend.core.extraction import ExtractionService
from backend.core.text_cleaning import TextCleaningService as TextCleaner
from backend.core.knowledge_core import KnowledgeCoreService
from backend.services.generators import ArtifactGenerator
from backend.services.pdf_renderer import PDFRenderer
from backend.models.artifacts import SlidesModel
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

load_dotenv()

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "frontend_mock", "public", "data")
os.makedirs(OUTPUT_DIR, exist_ok=True)

def render_pptx(slides_model: SlidesModel, filename: str):
    """
    Renders a SlidesModel to a PPTX file using python-pptx.
    """
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = slides_model.title
    subtitle.text = f"Audience: {slides_model.audience_level}"
    
    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_data in slides_model.slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        # Title
        title_shape = shapes.title
        title_shape.text = slide_data.heading
        
        # Content
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = slide_data.main_idea
        
        for point in slide_data.bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 1
            
        # Notes
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.speaker_notes

    out_path = os.path.join(OUTPUT_DIR, filename)
    prs.save(out_path)
    logger.info(f"Generated PPTX: {out_path}")

def main():
    logger.info("============================================================")
    logger.info("STEEL THREAD TEST: Real World Pipeline Verification")
    logger.info("============================================================")

    # 1. Initialize Services
    ingestor = IngestionService()
    extractor = ExtractionService()
    cleaner = TextCleaner()
    core_service = KnowledgeCoreService()
    generator = ArtifactGenerator()
    pdf_renderer = PDFRenderer(output_dir=OUTPUT_DIR)

    # 2. Ingest (YouTube)
    logger.info("STEP 1: Ingesting YouTube Video...")
    # NOTE: Using a shorter video for testing if needed, or the one requested.
    # User requested: https://www.youtube.com/watch?v=R9OCA6UFE-0 (Stoicism)
    url = "https://www.youtube.com/watch?v=R9OCA6UFE-0"
    user_id = "test_user"
    
    # ingest_meta = ingestor.process_youtube(url, user_id)
    # if ingest_meta.get("status") == "FAILED":
    #     logger.error("Ingest FAILED.")
    #     return
        
    # r2_key = ingest_meta.get("fileURL")
    # logger.info(f"STEP 1: SUCCESS. R2 Key: {r2_key}")
    
    # BYPASS INGEST (YouTube 403) - Use existing uploaded file
    r2_key = "uploads/dbf274d5-0648-4dff-9a3f-474b8fb58c13.wav" 
    logger.info(f"STEP 1: SKIPPED (Using existing Key: {r2_key})")

    # 3. Extract (Azure)
    logger.info("STEP 2: Extracting from R2 (Azure Transcription)...")
    text, ext_meta = extractor.extract_from_r2(r2_key)
    if len(text) < 50:
        logger.error(f"Extraction failed or too short. Text: {text}")
        return
        
    logger.info(f"STEP 2: SUCCESS. Extracted {len(text)} chars.")

    # 4. Clean (Gemini)
    logger.info("STEP 3: Cleaning Text...")
    cleaned_text = cleaner.clean_with_gemini(text)
    logger.info(f"STEP 3: SUCCESS. Cleaned text length: {len(cleaned_text)}")

    # 5. Knowledge Core (Gemini)
    logger.info("STEP 4: Generating Knowledge Core...")
    core = core_service.generate_knowledge_core(cleaned_text)
    if not core:
        logger.error("Knowledge Core generation FAILED.")
        return

    # Save Core for debugging
    core_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "knowledge_core.json")
    with open(core_path, "w") as f:
        f.write(core.model_dump_json(indent=2))
        
    logger.info(f"STEP 4: SUCCESS. Saved Knowledge Core to {core_path}")

    # 6. Generate & Render (The Big 5)
    logger.info("STEP 5: Generating Artifacts & Rendering...")

    # A. Final Exam -> PDF
    logger.info("  Generating Final Exam...")
    exam = generator.generate_exam(core)
    if exam:
        pdf_renderer.render_exam(exam, filename="final_exam")
        logger.info("  - Final Exam PDF: SUCCESS")
    else:
        logger.error("  - Final Exam Generation FAILED")

    # B. Quiz -> JSON
    logger.info("  Generating Quiz...")
    quiz = generator.generate_quiz(core)
    if quiz:
        with open(os.path.join(OUTPUT_DIR, "quiz.json"), "w") as f:
            f.write(quiz.model_dump_json(indent=2))
        logger.info("  - Quiz JSON: SUCCESS")
    else:
        logger.error("  - Quiz Generation FAILED")

    # C. Flashcards -> JSON
    logger.info("  Generating Flashcards...")
    cards = generator.generate_flashcards(core)
    if cards:
        with open(os.path.join(OUTPUT_DIR, "flashcards.json"), "w") as f:
            f.write(cards.model_dump_json(indent=2))
        logger.info("  - Flashcards JSON: SUCCESS")
    else:
        logger.error("  - Flashcards Generation FAILED")

    # D. Notes -> Markdown
    logger.info("  Generating Notes...")
    notes = generator.generate_notes(core)
    if notes:
        # Render simple markdown from NotesModel
        md_content = f"# {notes.title}\n\n"
        for section in notes.sections:
            md_content += f"## {section.heading}\n\n"
            md_content += "### Key Points\n"
            for kp in section.key_points:
                md_content += f"- {kp}\n"
            md_content += f"\n{section.content_block}\n\n"
        
        with open(os.path.join(OUTPUT_DIR, "notes.md"), "w") as f:
            f.write(md_content)
        # ALSO SAVE JSON for Frontend
        with open(os.path.join(OUTPUT_DIR, "notes.json"), "w") as f:
            f.write(notes.model_dump_json(indent=2))
        logger.info("  - Notes Markdown & JSON: SUCCESS")
    else:
        logger.error("  - Notes Generation FAILED")

    # E. Slides -> PPTX
    logger.info("  Generating Slides...")
    slides = generator.generate_slides(core)
    if slides:
        try:
            render_pptx(slides, "slides.pptx")
            # ALSO SAVE JSON for Frontend
            with open(os.path.join(OUTPUT_DIR, "slides.json"), "w") as f:
                f.write(slides.model_dump_json(indent=2))
            logger.info("  - Slides PPTX & JSON: SUCCESS")
        except Exception as e:
            logger.error(f"  - Slides Rendering FAILED: {e}")
    else:
        logger.error("  - Slides Generation FAILED")

    logger.info("============================================================")
    logger.info("STEEL THREAD COMPLETE. CHECK frontend_mock/public/data/")
    logger.info("============================================================")

if __name__ == "__main__":
    main()
