"""
BinaryRenderer: Handles rendering artifacts to binary formats (PDF, PPTX) and uploading to R2.

INVARIANTS:
- All binary outputs are uploaded to R2
- Storage paths follow pattern: {project_id}/artifacts/{artifact_id}.{ext}
- Original JSON data is preserved; binary is ADDED, not replaced
"""

import os
import uuid
import tempfile
import logging
from typing import Optional, Dict, Any

import boto3
from botocore.config import Config
from dotenv import load_dotenv

from backend.models.artifacts import FinalExamModel, SlidesModel
from backend.services.pdf_renderer import PDFRenderer

load_dotenv()

logger = logging.getLogger(__name__)


class BinaryRenderer:
    """
    Renders artifacts to binary formats and uploads to R2.
    
    Supported formats:
    - exam → PDF
    - slides → PPTX (future)
    """
    
    def __init__(self):
        self._setup_r2()
        self.pdf_renderer = PDFRenderer(output_dir=tempfile.gettempdir())
    
    def _setup_r2(self):
        """Initialize Cloudflare R2 client."""
        self.r2_endpoint = os.getenv("R2_ENDPOINT_URL")
        self.r2_key = os.getenv("R2_ACCESS_KEY_ID")
        self.r2_secret = os.getenv("R2_SECRET_ACCESS_KEY")
        self.r2_bucket = os.getenv("R2_BUCKET_NAME")

        if self.r2_endpoint and self.r2_key and self.r2_secret:
            self.s3_client = boto3.client(
                service_name='s3',
                endpoint_url=self.r2_endpoint,
                aws_access_key_id=self.r2_key,
                aws_secret_access_key=self.r2_secret,
                region_name='auto',
                config=Config(signature_version='s3v4')
            )
            self.r2_available = True
        else:
            logger.warning("R2 credentials missing. Binary uploads will be skipped.")
            self.s3_client = None
            self.r2_available = False
    
    def _upload_to_r2(self, file_path: str, object_key: str, content_type: str) -> str:
        """
        Uploads file to R2 with proper content type.
        Returns the R2 object key.
        """
        if not self.r2_available:
            raise RuntimeError("R2 not configured - cannot upload binary")
        
        logger.info(f"Uploading binary to R2: {object_key}")
        self.s3_client.upload_file(
            file_path, 
            self.r2_bucket, 
            object_key,
            ExtraArgs={'ContentType': content_type}
        )
        return object_key
    
    def render_exam_pdf(
        self, 
        exam: FinalExamModel, 
        project_id: uuid.UUID, 
        artifact_id: uuid.UUID
    ) -> Optional[Dict[str, Any]]:
        """
        Renders exam to PDF and uploads to R2.
        
        Returns:
            Binary metadata dict: { "format": "pdf", "storage_path": "...", "mime_type": "..." }
            or None if rendering fails or R2 unavailable
        """
        if not self.r2_available:
            logger.warning("R2 not available, skipping PDF generation")
            return None
        
        try:
            # Generate unique filename
            filename = f"exam_{artifact_id}"
            
            # Render to temp directory
            logger.info(f"Rendering exam PDF: {filename}")
            self.pdf_renderer.render_exam(exam, filename=filename)
            
            # Check if PDF was generated
            pdf_path = os.path.join(self.pdf_renderer.output_dir, f"{filename}.pdf")
            
            if not os.path.exists(pdf_path):
                # pdflatex might not be installed - check for .tex at least
                tex_path = os.path.join(self.pdf_renderer.output_dir, f"{filename}.tex")
                if os.path.exists(tex_path):
                    logger.warning("PDF not generated (pdflatex missing?), but .tex exists")
                    return None
                else:
                    logger.error("Neither PDF nor TEX file generated")
                    return None
            
            # Upload to R2
            object_key = f"{project_id}/artifacts/{artifact_id}.pdf"
            self._upload_to_r2(pdf_path, object_key, "application/pdf")
            
            # Clean up temp file
            try:
                os.remove(pdf_path)
                tex_path = os.path.join(self.pdf_renderer.output_dir, f"{filename}.tex")
                if os.path.exists(tex_path):
                    os.remove(tex_path)
            except Exception as e:
                logger.debug(f"Cleanup failed: {e}")
            
            logger.info(f"Exam PDF uploaded to R2: {object_key}")
            
            return {
                "format": "pdf",
                "storage_path": object_key,
                "mime_type": "application/pdf",
            }
            
        except Exception as e:
            logger.error(f"Failed to render exam PDF: {e}")
            return None
    
    def _validate_slide_schema(self, slides: SlidesModel) -> None:
        """
        Validate slide data before rendering.
        Raises ValueError with descriptive error if validation fails.
        """
        # Validate presentation-level fields
        if not slides.title or not isinstance(slides.title, str):
            raise ValueError(f"Slides missing valid title. Got: {type(slides.title).__name__} = {repr(slides.title)}")
        
        if not slides.slides or len(slides.slides) == 0:
            raise ValueError("Slides model has no slides array or it's empty")
        
        # Validate each slide
        for idx, slide in enumerate(slides.slides):
            # Heading is required and must be string
            if not hasattr(slide, 'heading') or not slide.heading:
                raise ValueError(f"Slide {idx}: missing 'heading' field")
            if not isinstance(slide.heading, str):
                raise ValueError(f"Slide {idx}: 'heading' must be str, got {type(slide.heading).__name__}")
            
            # bullet_points should be a list of strings (can be empty)
            if hasattr(slide, 'bullet_points') and slide.bullet_points:
                if not isinstance(slide.bullet_points, list):
                    raise ValueError(f"Slide {idx}: 'bullet_points' must be list, got {type(slide.bullet_points).__name__}")
                for bp_idx, bp in enumerate(slide.bullet_points):
                    if not isinstance(bp, str):
                        raise ValueError(f"Slide {idx}, bullet {bp_idx}: expected str, got {type(bp).__name__} = {repr(bp)}")
        
        logger.info(f"[PPTX] Schema validation passed: {len(slides.slides)} slides")

    def _render_smoke_test_pptx(self, project_id: uuid.UUID, artifact_id: uuid.UUID) -> Optional[Dict[str, Any]]:
        """
        Minimal smoke test: creates a 1-slide PPTX with hardcoded content.
        If this fails, the renderer infrastructure is broken.
        If this succeeds, the issue is with LLM slide content.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            logger.info("[PPTX SMOKE TEST] Starting hardcoded 1-slide test...")
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_box.text_frame.paragraphs[0].text = "Smoke Test Slide"
            title_box.text_frame.paragraphs[0].font.size = Pt(36)
            title_box.text_frame.paragraphs[0].font.bold = True
            
            # Body
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(4))
            body_box.text_frame.paragraphs[0].text = "• Renderer infrastructure OK"
            body_box.text_frame.paragraphs[0].font.size = Pt(24)
            
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                pptx_path = tmp.name
            
            prs.save(pptx_path)
            
            object_key = f"{project_id}/artifacts/{artifact_id}_smoketest.pptx"
            self._upload_to_r2(
                pptx_path,
                object_key,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
            os.remove(pptx_path)
            logger.info(f"[PPTX SMOKE TEST] SUCCESS - uploaded to {object_key}")
            return {"smoke_test": "passed", "path": object_key}
            
        except Exception as e:
            import traceback
            logger.error(f"[PPTX SMOKE TEST] FAILED: {e}")
            logger.error(f"[PPTX SMOKE TEST] Stack trace:\n{traceback.format_exc()}")
            return None

    def render_slides_pptx(
        self, 
        slides: SlidesModel, 
        project_id: uuid.UUID, 
        artifact_id: uuid.UUID
    ) -> Optional[Dict[str, Any]]:
        """
        Renders slides to PPTX and uploads to R2.
        
        Returns:
            Binary metadata dict: { "format": "pptx", "storage_path": "...", "mime_type": "..." }
            or None if rendering fails or R2 unavailable
        """
        import traceback
        
        if not self.r2_available:
            logger.warning("R2 not available, skipping PPTX generation")
            return None
        
        # --- Run smoke test first (can be disabled after debugging) ---
        smoke_result = self._render_smoke_test_pptx(project_id, artifact_id)
        if not smoke_result:
            logger.error("[PPTX] Smoke test failed - renderer infrastructure broken")
            return None
        logger.info("[PPTX] Smoke test passed - proceeding with real slides")
        
        # --- Validate schema before attempting render ---
        try:
            self._validate_slide_schema(slides)
        except ValueError as ve:
            logger.error(f"[PPTX] Schema validation failed: {ve}")
            return None
        
        # --- Log full slide data for debugging ---
        logger.info(f"[PPTX] Rendering {len(slides.slides)} slides for artifact {artifact_id}")
        for idx, s in enumerate(slides.slides):
            logger.debug(f"[PPTX] Slide {idx}: heading={repr(s.heading)}, bullets={len(s.bullet_points) if s.bullet_points else 0}")
        
        current_slide_idx = -1  # Track which slide we're on for error reporting
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            # Create presentation with fixed dimensions
            prs = Presentation()
            prs.slide_width = Inches(13.333)  # 16:9
            prs.slide_height = Inches(7.5)
            
            # Use blank layout for full control (like exam renderer)
            blank_layout = prs.slide_layouts[6]
            
            # Color scheme (ported from mock_frontend pptxgenjs)
            from pptx.dml.color import RGBColor
            COLOR_PRIMARY = RGBColor(0x1e, 0x1b, 0x4b)  # Deep indigo
            COLOR_ACCENT = RGBColor(0x63, 0x66, 0xf1)   # Indigo accent
            COLOR_TEXT = RGBColor(0x36, 0x36, 0x36)     # Dark gray
            COLOR_SUBTITLE = RGBColor(0xcb, 0xd5, 0xe1) # Light gray
            
            # --- Title Slide (ported from mock_frontend) ---
            current_slide_idx = 0
            title_slide = prs.slides.add_slide(blank_layout)
            
            # Background color
            background = title_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = COLOR_PRIMARY
            
            # Title text - centered at 30% height
            title_box = title_slide.shapes.add_textbox(Inches(0), Inches(2.25), Inches(13.333), Inches(1.5))
            tf = title_box.text_frame
            tf.word_wrap = True  # Enable word wrap
            p = tf.paragraphs[0]
            p.text = str(slides.title) if slides.title else "Untitled"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)  # White
            p.alignment = PP_ALIGN.CENTER
            
            # Subtitle - audience level at 50% height
            sub_box = title_slide.shapes.add_textbox(Inches(0), Inches(3.75), Inches(13.333), Inches(0.6))
            tf2 = sub_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = f"Audience: {str(slides.audience_level) if slides.audience_level else 'General'}"
            p2.font.size = Pt(20)
            p2.font.color.rgb = COLOR_SUBTITLE
            p2.alignment = PP_ALIGN.CENTER
            
            # Bottom accent bar
            from pptx.enum.shapes import MSO_SHAPE
            accent_bar = title_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = COLOR_ACCENT
            accent_bar.line.fill.background()
            
            # --- Content Slides (ported from mock_frontend pptxgenjs layout) ---
            for idx, slide_data in enumerate(slides.slides):
                current_slide_idx = idx + 1  # +1 because title slide is 0
                
                logger.debug(f"[PPTX] Rendering slide {current_slide_idx}: {repr(slide_data.heading)[:50]}")
                
                content_slide = prs.slides.add_slide(blank_layout)
                
                # Header bar background (light gray, like mock_frontend)
                header_bg = content_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)
                )
                header_bg.fill.solid()
                header_bg.fill.fore_color.rgb = RGBColor(0xf1, 0xf5, 0xf9)  # Light bg
                header_bg.line.fill.background()
                
                # Heading text - centered in header bar
                heading_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.8))
                h_tf = heading_box.text_frame
                h_tf.word_wrap = True  # Enable word wrap
                h_p = h_tf.paragraphs[0]
                h_p.text = str(slide_data.heading) if slide_data.heading else f"Slide {current_slide_idx}"
                h_p.font.size = Pt(32)
                h_p.font.bold = True
                h_p.font.color.rgb = COLOR_PRIMARY
                h_p.alignment = PP_ALIGN.CENTER
                
                # Accent line below header
                accent_line = content_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.1), Inches(11.333), Inches(0.03)
                )
                accent_line.fill.solid()
                accent_line.fill.fore_color.rgb = COLOR_ACCENT
                accent_line.line.fill.background()
                
                # Main idea - centered, italic (like mock_frontend)
                if hasattr(slide_data, 'main_idea') and slide_data.main_idea:
                    idea_box = content_slide.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10.333), Inches(0.8))
                    idea_tf = idea_box.text_frame
                    idea_tf.word_wrap = True  # Enable word wrap
                    idea_p = idea_tf.paragraphs[0]
                    idea_p.text = str(slide_data.main_idea)
                    idea_p.font.size = Pt(20)
                    idea_p.font.italic = True
                    idea_p.font.color.rgb = RGBColor(0x43, 0x38, 0xCA)  # Indigo
                    idea_p.alignment = PP_ALIGN.CENTER
                
                # Bullet points - centered block, left-aligned text (like mock_frontend)
                if hasattr(slide_data, 'bullet_points') and slide_data.bullet_points:
                    # Wider margins (1" each side) to prevent text overflow
                    bullet_box = content_slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.333), Inches(4.5))
                    b_tf = bullet_box.text_frame
                    b_tf.word_wrap = True  # Critical for text wrapping
                    
                    for bp_idx, bp in enumerate(slide_data.bullet_points):
                        bp_text = str(bp) if bp else ""
                        
                        if bp_idx == 0:
                            p = b_tf.paragraphs[0]
                        else:
                            p = b_tf.add_paragraph()
                        
                        p.text = f"{bp_idx + 1}. {bp_text}"  # Numbered like mock_frontend
                        p.font.size = Pt(20)
                        p.font.color.rgb = COLOR_TEXT
                        p.space_after = Pt(16)
                
                # Speaker notes
                if hasattr(slide_data, 'speaker_notes') and slide_data.speaker_notes:
                    try:
                        notes_slide = content_slide.notes_slide
                        notes_slide.notes_text_frame.text = str(slide_data.speaker_notes)
                    except Exception as notes_err:
                        logger.warning(f"[PPTX] Failed to add speaker notes to slide {current_slide_idx}: {notes_err}")
            
            # --- Save to temp file ---
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                pptx_path = tmp.name
            
            logger.info(f"[PPTX] Saving to temp file: {pptx_path}")
            prs.save(pptx_path)
            
            # --- Upload to R2 ---
            object_key = f"{project_id}/artifacts/{artifact_id}.pptx"
            logger.info(f"[PPTX] Uploading to R2: {object_key}")
            self._upload_to_r2(
                pptx_path, 
                object_key, 
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
            # Cleanup temp file
            try:
                os.remove(pptx_path)
            except Exception as cleanup_err:
                logger.debug(f"[PPTX] Cleanup failed: {cleanup_err}")
            
            logger.info(f"[PPTX] SUCCESS - uploaded {len(slides.slides) + 1} slides to {object_key}")
            
            return {
                "format": "pptx",
                "storage_path": object_key,
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
            
        except ImportError as ie:
            logger.error(f"[PPTX] ImportError: {ie}")
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            return None
            
        except Exception as e:
            # --- Full error logging ---
            logger.error(f"[PPTX] FAILED at slide index {current_slide_idx}")
            logger.error(f"[PPTX] Exception type: {type(e).__name__}")
            logger.error(f"[PPTX] Exception message: {e}")
            logger.error(f"[PPTX] Full stack trace:\n{traceback.format_exc()}")
            
            # Log the problematic slide data
            if current_slide_idx > 0 and current_slide_idx <= len(slides.slides):
                bad_slide = slides.slides[current_slide_idx - 1]  # -1 because title is 0
                logger.error(f"[PPTX] Failing slide data:")
                logger.error(f"  heading: {repr(bad_slide.heading)}")
                logger.error(f"  main_idea: {repr(getattr(bad_slide, 'main_idea', None))}")
                logger.error(f"  bullet_points: {repr(getattr(bad_slide, 'bullet_points', None))}")
                logger.error(f"  speaker_notes: {repr(getattr(bad_slide, 'speaker_notes', None))[:100]}...")
            
            return None
